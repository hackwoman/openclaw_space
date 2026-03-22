#!/usr/bin/env python3
"""Office Assistant Demo - 展示办公秘书完整能力"""

import json
import os
from datetime import datetime

output_dir = '/home/kui/.openclaw/workspace/office_assistant_output'
os.makedirs(output_dir, exist_ok=True)

print("🏢 办公秘书能力演示")
print("=" * 60)

# ============ 1. 工作汇报 ============
print("\n📝 1. 创建工作汇报...")

report_data = {
    "title": "2025年第一季度销售部工作汇报",
    "author": "张三",
    "date": "2025年4月1日",
    "overview": "本季度销售部在公司领导的正确指导下，紧紧围绕年度销售目标，积极开拓市场，优化销售策略，取得了较好的业绩。现将本季度工作情况汇报如下：",
    "achievements": [
        "完成销售额1,280万元，完成季度目标的115%，同比增长23%",
        "新开发客户56家，其中重点客户12家",
        "客户满意度达到94.5%，较上季度提升2.3个百分点",
        "成功签约3个百万级大客户合同",
        "销售团队人均业绩提升18%"
    ],
    "data": {
        "headers": ["指标", "目标值", "实际值", "完成率"],
        "rows": [
            ["销售额", "1,100万", "1,280万", "116%"],
            ["新客户数", "40家", "56家", "140%"],
            ["回款率", "90%", "92%", "102%"],
            ["客户满意度", "90%", "94.5%", "105%"]
        ]
    },
    "issues": [
        "部分区域市场开拓进度滞后，需加强资源投入",
        "团队新人较多，培训周期较长",
        "竞争对手降价压力增大"
    ],
    "plans": [
        "加大重点区域市场开拓力度",
        "优化新人培训体系，缩短成长周期",
        "推出差异化产品组合，提升竞争力",
        "加强客户关系维护，提高复购率"
    ]
}

# Create Word document
import sys
sys.path.insert(0, '/home/kui/.openclaw/workspace/skills/office-assistant/scripts')
from create_document import create_report

create_report(report_data, f'{output_dir}/工作汇报.docx')

# ============ 2. 项目方案 ============
print("\n📋 2. 创建项目方案...")

proposal_data = {
    "title": "企业数字化转型项目方案",
    "background": "随着数字经济的快速发展，我司现有的信息化系统已难以满足业务增长需求。为提升运营效率、降低成本、增强竞争力，特制定本数字化转型方案。",
    "objectives": [
        "实现核心业务流程数字化，提升运营效率30%以上",
        "建立统一的数据平台，实现数据驱动决策",
        "提升客户体验，客户满意度达到95%以上",
        "降低运营成本15%以上"
    ],
    "scope": "本项目涵盖财务管理系统、客户关系管理（CRM）、供应链管理（SCM）、人力资源管理（HRM）四大核心系统的建设与集成。",
    "timeline": [
        {"phase": "需求调研", "time": "第1-2周", "work": "业务需求收集、流程梳理"},
        {"phase": "系统设计", "time": "第3-4周", "work": "架构设计、技术选型"},
        {"phase": "开发实施", "time": "第5-12周", "work": "系统开发、单元测试"},
        {"phase": "系统集成", "time": "第13-14周", "work": "系统对接、集成测试"},
        {"phase": "试运行", "time": "第15-16周", "work": "用户培训、试运行"},
        {"phase": "正式上线", "time": "第17周", "work": "全面上线、运维支持"}
    ],
    "budget": [
        {"item": "软件开发", "amount": "800,000", "note": "含人力成本"},
        {"item": "硬件设备", "amount": "300,000", "note": "服务器、存储等"},
        {"item": "软件许可", "amount": "200,000", "note": "第三方软件"},
        {"item": "培训费用", "amount": "50,000", "note": "用户培训"},
        {"item": "其他费用", "amount": "50,000", "note": "差旅、杂费"}
    ],
    "risks": [
        {"risk": "需求变更频繁", "mitigation": "建立需求变更管理流程"},
        {"risk": "技术难度超预期", "mitigation": "提前进行技术验证"},
        {"risk": "人员投入不足", "mitigation": "制定详细人员计划"}
    ],
    "conclusion": "本项目的实施将显著提升我司的数字化水平，为业务发展提供有力支撑。建议尽快启动项目，抢占市场先机。"
}

from create_document import create_proposal
create_proposal(proposal_data, f'{output_dir}/项目方案.docx')

# ============ 3. 会议纪要 ============
print("\n📅 3. 创建会议纪要...")

minutes_data = {
    "title": "2025年第一季度经营分析会",
    "time": "2025年4月5日 14:00-17:00",
    "location": "公司三楼会议室",
    "chairman": "李总",
    "attendees": ["王副总", "张经理", "刘经理", "陈经理", "赵经理", "周秘书"],
    "agenda": ["第一季度经营情况回顾", "各部门工作汇报", "存在问题分析", "第二季度工作部署"],
    "discussion": [
        {"topic": "销售部工作", "content": "本季度销售额1,280万，超额完成目标。重点客户开发取得突破。"},
        {"topic": "生产部工作", "content": "产能利用率92%，产品质量合格率99.5%。设备维护计划执行良好。"},
        {"topic": "财务状况", "content": "毛利率38.5%，净利润率12.3%。现金流状况良好。"}
    ],
    "decisions": [
        "同意销售部增加市场推广预算20%",
        "批准生产部设备升级计划",
        "要求各部门在4月15日前提交第二季度详细计划"
    ],
    "actions": [
        {"task": "完成市场推广方案", "owner": "张经理", "deadline": "4月10日", "status": "进行中"},
        {"task": "提交设备采购申请", "owner": "刘经理", "deadline": "4月12日", "status": "待开始"},
        {"task": "制定Q2部门计划", "owner": "各部门负责人", "deadline": "4月15日", "status": "进行中"}
    ]
}

from create_document import create_meeting_minutes
create_meeting_minutes(minutes_data, f'{output_dir}/会议纪要.docx')

# ============ 4. 通知公文 ============
print("\n📢 4. 创建通知公文...")

memo_data = {
    "company": "XX科技有限公司",
    "doc_type": "通知",
    "year": "2025",
    "number": "015",
    "title": "关于举办2025年度员工培训的通知",
    "content": "为提升员工专业技能和综合素质，促进公司持续发展，经研究决定，举办2025年度员工培训活动。现将有关事项通知如下：\n\n一、培训对象：公司全体员工\n二、培训时间：2025年4月20日-5月20日\n三、培训内容：专业技能培训、管理能力提升、企业文化培训\n四、培训方式：线上+线下相结合\n\n请各部门合理安排工作，确保员工按时参加培训。",
    "recipients": ["各部门", "全体员工"],
    "sender": "人力资源部",
    "date": "2025年4月8日"
}

from create_document import create_memo
create_memo(memo_data, f'{output_dir}/通知.docx')

# ============ 5. 数据分析报告 ============
print("\n📊 5. 创建数据分析报告...")

import pandas as pd
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import seaborn as sns

# Create sample data
data = {
    '月份': ['1月', '2月', '3月', '4月', '5月', '6月'],
    '销售额': [180, 150, 220, 280, 310, 350],
    '成本': [100, 90, 120, 150, 170, 190],
    '利润': [80, 60, 100, 130, 140, 160],
    '客户数': [45, 38, 52, 68, 75, 88]
}
df = pd.DataFrame(data)

# Create charts
fig, axes = plt.subplots(2, 2, figsize=(14, 10))

# Sales trend
axes[0, 0].plot(df['月份'], df['销售额'], marker='o', linewidth=2, color='#4a90d9')
axes[0, 0].fill_between(df['月份'], df['销售额'], alpha=0.3, color='#4a90d9')
axes[0, 0].set_title('销售趋势', fontsize=14, fontweight='bold')
axes[0, 0].set_ylabel('万元')

# Revenue vs Cost
x = range(len(df['月份']))
width = 0.35
axes[0, 1].bar([i - width/2 for i in x], df['销售额'], width, label='销售额', color='#4a90d9')
axes[0, 1].bar([i + width/2 for i in x], df['成本'], width, label='成本', color='#ff6b6b')
axes[0, 1].set_xticks(x)
axes[0, 1].set_xticklabels(df['月份'])
axes[0, 1].set_title('收入与成本对比', fontsize=14, fontweight='bold')
axes[0, 1].legend()

# Profit trend
axes[1, 0].bar(df['月份'], df['利润'], color='#00b894')
axes[1, 0].set_title('利润趋势', fontsize=14, fontweight='bold')
axes[1, 0].set_ylabel('万元')

# Customer growth
axes[1, 1].plot(df['月份'], df['客户数'], marker='s', linewidth=2, color='#a29bfe', markersize=8)
axes[1, 1].fill_between(df['月份'], df['客户数'], alpha=0.3, color='#a29bfe')
axes[1, 1].set_title('客户增长', fontsize=14, fontweight='bold')
axes[1, 1].set_ylabel('客户数')

plt.suptitle('2025年上半年经营数据分析', fontsize=16, fontweight='bold')
plt.tight_layout()
plt.savefig(f'{output_dir}/数据分析图表.png', dpi=150, bbox_inches='tight')
plt.close()
print(f"✅ 数据分析图表已保存: 数据分析图表.png")

# Save data to Excel
with pd.ExcelWriter(f'{output_dir}/数据分析.xlsx', engine='openpyxl') as writer:
    df.to_excel(writer, sheet_name='月度数据', index=False)
    
    # Summary stats
    summary = df.describe()
    summary.to_excel(writer, sheet_name='统计摘要')
print(f"✅ 数据分析已保存: 数据分析.xlsx")

# ============ 6. PPT演示文稿 ============
print("\n📊 6. 创建PPT演示文稿...")

from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB

prs = Presentation()

# Slide 1: Title
slide1 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide1.background.fill
bg.solid()
bg.fore_color.rgb = PptRGB(26, 26, 46)

box = slide1.shapes.add_textbox(PptInches(1), PptInches(2), PptInches(8), PptInches(1.5))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "2025年第一季度"
p.font.size = PptPt(36)
p.font.color.rgb = PptRGB(255, 215, 0)
p.alignment = 1

box = slide1.shapes.add_textbox(PptInches(1), PptInches(3.5), PptInches(8), PptInches(1))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "经营分析报告"
p.font.size = PptPt(28)
p.font.color.rgb = PptRGB(255, 255, 255)
p.alignment = 1

# Slide 2: Key Metrics
slide2 = prs.slides.add_slide(prs.slide_layouts[6])
box = slide2.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(0.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "📊 核心指标"
p.font.size = PptPt(28)
p.font.bold = True
p.font.color.rgb = PptRGB(255, 107, 107)

metrics = [
    ("1,280万", "销售额", "0, 180, 148"),
    ("116%", "目标完成率", "100, 100, 255"),
    ("56家", "新客户数", "255, 215, 0"),
    ("94.5%", "客户满意度", "162, 155, 254"),
]

for i, (value, label, color) in enumerate(metrics):
    x = 0.5 + (i % 2) * 4.5
    y = 1.5 + (i // 2) * 2.5
    
    shape = slide2.shapes.add_shape(1, PptInches(x), PptInches(y), PptInches(4), PptInches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = PptRGB(45, 45, 60)
    shape.line.fill.background()
    
    box = slide2.shapes.add_textbox(PptInches(x + 0.3), PptInches(y + 0.3), PptInches(3.4), PptInches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = PptPt(32)
    p.font.bold = True
    r, g, b = map(int, color.split(', '))
    p.font.color.rgb = PptRGB(r, g, b)
    p.alignment = 1
    
    box = slide2.shapes.add_textbox(PptInches(x + 0.3), PptInches(y + 1.3), PptInches(3.4), PptInches(0.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = label
    p.font.size = PptPt(14)
    p.font.color.rgb = PptRGB(200, 200, 200)
    p.alignment = 1

# Slide 3: Summary
slide3 = prs.slides.add_slide(prs.slide_layouts[6])
box = slide3.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(0.8))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "✅ 工作总结"
p.font.size = PptPt(28)
p.font.bold = True
p.font.color.rgb = PptRGB(0, 206, 201)

summary_items = [
    "销售额超额完成，同比增长23%",
    "新客户开发成效显著，超额40%",
    "客户满意度持续提升",
    "团队建设取得进展"
]

for i, item in enumerate(summary_items):
    box = slide3.shapes.add_textbox(PptInches(1), PptInches(1.5 + i * 0.8), PptInches(8), PptInches(0.6))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = f"✓ {item}"
    p.font.size = PptPt(18)
    p.font.color.rgb = PptRGB(220, 220, 220)

# Slide 4: Thank you
slide4 = prs.slides.add_slide(prs.slide_layouts[6])
bg = slide4.background.fill
bg.solid()
bg.fore_color.rgb = PptRGB(26, 26, 46)

box = slide4.shapes.add_textbox(PptInches(1), PptInches(2.5), PptInches(8), PptInches(1.5))
tf = box.text_frame
p = tf.paragraphs[0]
p.text = "感谢聆听！"
p.font.size = PptPt(40)
p.font.bold = True
p.font.color.rgb = PptRGB(255, 215, 0)
p.alignment = 1

prs.save(f'{output_dir}/经营分析汇报.pptx')
print(f"✅ PPT演示已保存: 经营分析汇报.pptx")

# ============ Summary ============
print("\n" + "=" * 60)
print("🎉 办公秘书演示完成！")
print(f"\n📁 所有文件保存在: {output_dir}")
print("\n生成的文件：")
for f in sorted(os.listdir(output_dir)):
    size = os.path.getsize(f'{output_dir}/{f}')
    print(f"  📄 {f} ({size/1024:.1f} KB)")
