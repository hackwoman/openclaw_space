#!/usr/bin/env python3
"""Office 文档处理能力展示"""

from openpyxl import Workbook
from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
from openpyxl.chart import BarChart, Reference
from docx import Document
from docx.shared import Inches, Pt, RGBColor
from docx.enum.text import WD_ALIGN_PARAGRAPH
from pptx import Presentation
from pptx.util import Inches as PptInches, Pt as PptPt
from pptx.dml.color import RGBColor as PptRGB
import os

output_dir = '/home/kui/.openclaw/workspace/office_files'
os.makedirs(output_dir, exist_ok=True)

# ============ 1. Excel 演示 ============
def create_excel_demo():
    wb = Workbook()
    
    # ---- 销售数据表 ----
    ws = wb.active
    ws.title = "销售数据"
    
    # 表头样式
    header_font = Font(bold=True, color="FFFFFF", size=12)
    header_fill = PatternFill(start_color="4472C4", end_color="4472C4", fill_type="solid")
    thin_border = Border(
        left=Side(style='thin'),
        right=Side(style='thin'),
        top=Side(style='thin'),
        bottom=Side(style='thin')
    )
    
    # 表头
    headers = ['产品', '1月', '2月', '3月', '合计', '增长率']
    for col, header in enumerate(headers, 1):
        cell = ws.cell(row=1, column=col, value=header)
        cell.font = header_font
        cell.fill = header_fill
        cell.alignment = Alignment(horizontal='center')
        cell.border = thin_border
    
    # 数据
    products = [
        ['笔记本电脑', 150, 180, 200],
        ['智能手机', 280, 320, 350],
        ['平板电脑', 90, 110, 130],
        ['智能手表', 60, 75, 95],
        ['耳机', 200, 240, 280],
    ]
    
    for row, product in enumerate(products, 2):
        for col, value in enumerate(product, 1):
            cell = ws.cell(row=row, column=col, value=value)
            cell.border = thin_border
            cell.alignment = Alignment(horizontal='center')
        
        # 合计公式
        ws.cell(row=row, column=5).value = f'=SUM(B{row}:D{row})'
        ws.cell(row=row, column=5).border = thin_border
        ws.cell(row=row, column=5).alignment = Alignment(horizontal='center')
        
        # 增长率
        ws.cell(row=row, column=6).value = f'=(D{row}-B{row})/B{row}'
        ws.cell(row=row, column=6).number_format = '0.0%'
        ws.cell(row=row, column=6).border = thin_border
        ws.cell(row=row, column=6).alignment = Alignment(horizontal='center')
    
    # 添加图表
    chart = BarChart()
    chart.type = "col"
    chart.title = "产品销售对比"
    data = Reference(ws, min_col=2, max_col=4, min_row=1, max_row=6)
    categories = Reference(ws, min_col=1, min_row=2, max_row=6)
    chart.add_data(data, titles_from_data=True)
    chart.set_categories(categories)
    chart.shape = 4
    ws.add_chart(chart, "H2")
    
    # 调整列宽
    ws.column_dimensions['A'].width = 15
    for col in ['B', 'C', 'D', 'E', 'F']:
        ws.column_dimensions[col].width = 12
    
    # ---- 数据透视表 ----
    ws2 = wb.create_sheet("统计分析")
    ws2['A1'] = "📊 销售统计分析"
    ws2['A1'].font = Font(bold=True, size=16, color="4472C4")
    ws2.merge_cells('A1:D1')
    
    stats = [
        ['指标', '数值'],
        ['总销售额', '=销售数据!E7'],
        ['月均销售', '=AVERAGE(销售数据!B2:D6)'],
        ['最高单品', '智能手机'],
        ['增长率最高', '智能手表'],
    ]
    
    for row, data in enumerate(stats, 3):
        for col, value in enumerate(data, 1):
            cell = ws2.cell(row=row, column=col, value=value)
            cell.border = thin_border
            if row == 3:
                cell.font = header_font
                cell.fill = header_fill
    
    ws2.column_dimensions['A'].width = 15
    ws2.column_dimensions['B'].width = 15
    
    wb.save(f'{output_dir}/销售报告.xlsx')
    print("✅ Excel 文件已保存: 销售报告.xlsx")

# ============ 2. Word 演示 ============
def create_word_demo():
    doc = Document()
    
    # 标题
    title = doc.add_heading('📋 项目报告', 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    
    # 副标题
    subtitle = doc.add_paragraph()
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    run = subtitle.add_run('OpenClaw AI 助手 - 2026年3月')
    run.font.size = Pt(14)
    run.font.color.rgb = RGBColor(100, 100, 100)
    
    doc.add_paragraph()  # 空行
    
    # 项目概述
    doc.add_heading('一、项目概述', level=1)
    doc.add_paragraph(
        '本项目旨在开发一个智能化的 AI 助手系统，'
        '能够帮助用户处理日常任务、回答问题、'
        '并提供各种工具和技能支持。'
    )
    
    # 主要功能
    doc.add_heading('二、主要功能', level=1)
    
    features = [
        ('🎨 图像生成', '支持 SVG、PNG 格式，可创建卡通、像素艺术、图表等'),
        ('📊 数据处理', '支持 Excel、Word、PPT 等 Office 文档处理'),
        ('🔍 信息搜索', '支持网页搜索和内容提取'),
        ('🤖 智能对话', '基于大语言模型的自然语言交互'),
    ]
    
    for title, desc in features:
        p = doc.add_paragraph()
        run = p.add_run(f'✓ {title}')
        run.bold = True
        run.font.color.rgb = RGBColor(0, 180, 148)
        p.add_run(f'\n  {desc}')
    
    # 开发进度
    doc.add_heading('三、开发进度', level=1)
    
    table = doc.add_table(rows=5, cols=3)
    table.style = 'Light Grid Accent 1'
    
    headers = ['阶段', '状态', '完成度']
    for i, header in enumerate(headers):
        table.rows[0].cells[i].text = header
    
    progress = [
        ['需求分析', '✅ 已完成', '100%'],
        ['系统设计', '✅ 已完成', '100%'],
        ['核心开发', '🔄 进行中', '85%'],
        ['测试优化', '📅 计划中', '30%'],
    ]
    
    for row_idx, row_data in enumerate(progress, 1):
        for col_idx, cell_data in enumerate(row_data):
            table.rows[row_idx].cells[col_idx].text = cell_data
    
    # 总结
    doc.add_heading('四、总结与展望', level=1)
    doc.add_paragraph(
        '目前项目进展顺利，核心功能已基本实现。'
        '下一步将重点优化用户体验，增加更多实用技能，'
        '并持续提升系统的稳定性和响应速度。'
    )
    
    # 签名
    doc.add_paragraph()
    p = doc.add_paragraph()
    p.alignment = WD_ALIGN_PARAGRAPH.RIGHT
    run = p.add_run('报告人：OpenClaw AI')
    run.font.color.rgb = RGBColor(100, 100, 100)
    
    doc.save(f'{output_dir}/项目报告.docx')
    print("✅ Word 文件已保存: 项目报告.docx")

# ============ 3. PPT 演示 ============
def create_ppt_demo():
    prs = Presentation()
    
    # ---- 封面页 ----
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    
    # 背景色
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PptRGB(26, 26, 46)
    
    # 标题
    title_box = slide.shapes.add_textbox(PptInches(1), PptInches(1.5), PptInches(8), PptInches(1.5))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🦞 OpenClaw AI 助手"
    p.font.size = PptPt(44)
    p.font.bold = True
    p.font.color.rgb = PptRGB(255, 215, 0)
    p.alignment = 1  # 居中
    
    # 副标题
    sub_box = slide.shapes.add_textbox(PptInches(1), PptInches(3.2), PptInches(8), PptInches(1))
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.text = "智能 · 高效 · 全能"
    p.font.size = PptPt(24)
    p.font.color.rgb = PptRGB(255, 255, 255)
    p.alignment = 1
    
    # ---- 功能介绍页 ----
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 标题
    title_box = slide2.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "✨ 核心功能"
    p.font.size = PptPt(32)
    p.font.bold = True
    p.font.color.rgb = PptRGB(255, 107, 107)
    
    # 功能卡片
    features = [
        ("🎨 图像生成", "SVG/PNG/像素艺术", "0, 180, 148"),
        ("📊 数据处理", "Excel/Word/PPT", "100, 100, 255"),
        ("🔍 智能搜索", "网页搜索与分析", "255, 215, 0"),
        ("💬 智能对话", "自然语言交互", "162, 155, 254"),
    ]
    
    for i, (title, desc, color) in enumerate(features):
        x = 0.5 + (i % 2) * 4.7
        y = 1.5 + (i // 2) * 2.2
        
        # 卡片背景
        shape = slide2.shapes.add_shape(
            1,  # 矩形
            PptInches(x), PptInches(y),
            PptInches(4.2), PptInches(1.8)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = PptRGB(45, 45, 60)
        shape.line.fill.background()
        
        # 标题
        box = slide2.shapes.add_textbox(PptInches(x + 0.3), PptInches(y + 0.2), PptInches(3.6), PptInches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = title
        p.font.size = PptPt(20)
        p.font.bold = True
        r, g, b = map(int, color.split(', '))
        p.font.color.rgb = PptRGB(r, g, b)
        
        # 描述
        box = slide2.shapes.add_textbox(PptInches(x + 0.3), PptInches(y + 0.9), PptInches(3.6), PptInches(0.6))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = desc
        p.font.size = PptPt(14)
        p.font.color.rgb = PptRGB(200, 200, 200)
    
    # ---- 数据页 ----
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    
    title_box = slide3.shapes.add_textbox(PptInches(0.5), PptInches(0.3), PptInches(9), PptInches(0.8))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📊 使用统计"
    p.font.size = PptPt(32)
    p.font.bold = True
    p.font.color.rgb = PptRGB(0, 206, 201)
    
    # 统计数据
    stats = [
        ("10+", "支持格式", "255, 107, 107"),
        ("50+", "技能工具", "0, 206, 201"),
        ("100%", "免费使用", "255, 215, 0"),
        ("24/7", "全天在线", "162, 155, 254"),
    ]
    
    for i, (number, label, color) in enumerate(stats):
        x = 0.5 + i * 2.3
        
        box = slide3.shapes.add_textbox(PptInches(x), PptInches(1.8), PptInches(2), PptInches(1))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = number
        p.font.size = PptPt(48)
        p.font.bold = True
        r, g, b = map(int, color.split(', '))
        p.font.color.rgb = PptRGB(r, g, b)
        p.alignment = 1
        
        box = slide3.shapes.add_textbox(PptInches(x), PptInches(2.8), PptInches(2), PptInches(0.5))
        tf = box.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = PptPt(16)
        p.font.color.rgb = PptRGB(200, 200, 200)
        p.alignment = 1
    
    # ---- 结尾页 ----
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    
    background = slide4.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PptRGB(26, 26, 46)
    
    box = slide4.shapes.add_textbox(PptInches(1), PptInches(2), PptInches(8), PptInches(1.5))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "感谢使用 OpenClaw! 🦞"
    p.font.size = PptPt(36)
    p.font.bold = True
    p.font.color.rgb = PptRGB(255, 215, 0)
    p.alignment = 1
    
    box = slide4.shapes.add_textbox(PptInches(1), PptInches(3.5), PptInches(8), PptInches(1))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = "智能助手，无限可能"
    p.font.size = PptPt(20)
    p.font.color.rgb = PptRGB(200, 200, 200)
    p.alignment = 1
    
    prs.save(f'{output_dir}/产品介绍.pptx')
    print("✅ PPT 文件已保存: 产品介绍.pptx")

# 运行所有演示
if __name__ == '__main__':
    print("📄 开始生成 Office 演示文档...")
    print()
    create_excel_demo()
    create_word_demo()
    create_ppt_demo()
    print()
    print("🎉 所有 Office 文档生成完成！")
